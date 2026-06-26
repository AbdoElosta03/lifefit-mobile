# -*- coding: utf-8 -*-
"""Generate FitLife graduation presentation as a PowerPoint (.pptx).
Theme: FitLife teal | Font: Tajawal | Direction: RTL.
"""
import os
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BASE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(BASE, "images")

# ---------------- palette ----------------
TEAL900 = RGBColor(0x0B, 0x4F, 0x44)
TEAL800 = RGBColor(0x0F, 0x6B, 0x5C)
TEAL700 = RGBColor(0x0D, 0x94, 0x88)
TEAL600 = RGBColor(0x11, 0xA3, 0x94)
TEAL500 = RGBColor(0x14, 0xB8, 0xA6)
MINT100 = RGBColor(0xCC, 0xFB, 0xF1)
INK     = RGBColor(0x0F, 0x2A, 0x26)
BODY    = RGBColor(0x33, 0x41, 0x55)
MUTED   = RGBColor(0x6B, 0x7C, 0x86)
LINE    = RGBColor(0xDD, 0xE7, 0xE4)
BG      = RGBColor(0xF4, 0xF8, 0xF7)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
ROSE    = RGBColor(0xEF, 0x44, 0x44)
BLUE    = RGBColor(0x3B, 0x82, 0xF6)
VIOLET  = RGBColor(0x8B, 0x5C, 0xF6)

TINTS = {
    'teal':   (RGBColor(0xD7, 0xF7, 0xEF), TEAL700),
    'amber':  (RGBColor(0xFE, 0xF3, 0xD7), AMBER),
    'rose':   (RGBColor(0xFD, 0xE0, 0xE0), ROSE),
    'blue':   (RGBColor(0xDD, 0xEA, 0xFE), BLUE),
    'violet': (RGBColor(0xEA, 0xE3, 0xFD), VIOLET),
}

FONT   = "Tajawal"
FONT_B = "Tajawal"        # bold handled via run.bold
EMOJI  = "Segoe UI Emoji"

PW, PH = 13.333, 7.5
MX = 0.62

prs = Presentation()
prs.slide_width = Inches(PW)
prs.slide_height = Inches(PH)
BLANK = prs.slide_layouts[6]

# ---------------- helpers ----------------

def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s

def gradient(shape, c1, c2, angle=45):
    shape.fill.gradient()
    stops = shape.fill.gradient_stops
    stops[0].position = 0.0; stops[0].color.rgb = c1
    stops[1].position = 1.0; stops[1].color.rgb = c2
    try:
        shape.fill.gradient_angle = angle
    except Exception:
        pass
    shape.line.fill.background()
    shape.shadow.inherit = False

def no_line(shape):
    shape.line.fill.background()

def soft_shadow(shape):
    """Apply a subtle outer shadow using a schema-valid effectLst position."""
    spPr = shape._element.spPr
    # effectLst must come after <a:ln>; append at end is correct for our shapes
    existing = spPr.find(qn('a:effectLst'))
    if existing is not None:
        return
    el = etree.SubElement(spPr, qn('a:effectLst'))
    sh = etree.SubElement(el, qn('a:outerShdw'),
                          {'blurRad': '120000', 'dist': '40000', 'dir': '5400000', 'rotWithShape': '0'})
    clr = etree.SubElement(sh, qn('a:srgbClr'), {'val': '0D5046'})
    etree.SubElement(clr, qn('a:alpha'), {'val': '20000'})

def _cs(run, name):
    rPr = run.font._element  # CT_TextCharacterProperties (the rPr itself)
    cs = rPr.find(qn('a:cs'))
    if cs is None:
        cs = etree.SubElement(rPr, qn('a:cs'))
    cs.set('typeface', name)

def set_rtl(p):
    pPr = p._p.get_or_add_pPr()
    pPr.set('rtl', '1')

def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return box, tf

def para(tf, runs, size=18, color=BODY, bold=False, align=PP_ALIGN.RIGHT,
         rtl=True, first=False, space_before=0, space_after=4, line=1.12, font=FONT):
    """runs: str OR list of (text, color|None, bold|None, font|None)."""
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    if rtl:
        set_rtl(p)
    p.space_before = Pt(space_before); p.space_after = Pt(space_after)
    try:
        p.line_spacing = line
    except Exception:
        pass
    if isinstance(runs, str):
        runs = [(runs, color, bold, font)]
    for item in runs:
        text = item[0]
        rc = item[1] if len(item) > 1 and item[1] is not None else color
        rb = item[2] if len(item) > 2 and item[2] is not None else bold
        rf = item[3] if len(item) > 3 and item[3] is not None else font
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = rb
        run.font.name = rf; run.font.color.rgb = rc
        _cs(run, rf)
    return p

def rrect(s, x, y, w, h, fill=WHITE, line=LINE, radius=0.08, line_w=1.0, shadow=False):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    if shadow:
        soft_shadow(sh)
    return sh

def icon(s, x, y, emoji, tint='teal', d=0.66):
    bgc, fgc = TINTS[tint]
    sq = rrect(s, x, y, d, d, fill=bgc, line=None, radius=0.28)
    _, tf = tb(s, x, y, d, d, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(emoji, fgc, False, EMOJI)], size=22, align=PP_ALIGN.CENTER, rtl=False, first=True, space_after=0)
    return sq

def header(s, num, kicker, title_runs, light=False):
    badge_w = 0.92
    bx = PW - MX - badge_w
    b = rrect(s, bx, 0.5, badge_w, badge_w, radius=0.22)
    if light:
        b.fill.solid(); b.fill.fore_color.rgb = WHITE
        b.line.color.rgb = WHITE
    else:
        gradient(b, TEAL500, TEAL800, 45)
    _, tf = tb(s, bx, 0.5, badge_w, badge_w, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(num, (TEAL800 if light else WHITE), True, FONT)], size=24,
         align=PP_ALIGN.CENTER, rtl=False, first=True, space_after=0)

    tx = MX
    tw = bx - 0.22 - MX
    _, tf2 = tb(s, tx, 0.5, tw, 1.05)
    kcolor = WHITE if light else TEAL700
    para(tf2, [(kicker, kcolor, True, FONT)], size=14, align=PP_ALIGN.RIGHT, first=True, space_after=2)
    para(tf2, title_runs, size=28, bold=True, align=PP_ALIGN.RIGHT, space_after=0)

def footer(s, light=False):
    c = RGBColor(0xCF, 0xEF, 0xE8) if light else MUTED
    sq = rrect(s, MX, PH - 0.5, 0.16, 0.16, fill=(WHITE if light else TEAL700), line=None, radius=0.25)
    _, tf = tb(s, MX + 0.24, PH - 0.56, 2.0, 0.3)
    para(tf, [("FitLife", c, True, FONT)], size=12, align=PP_ALIGN.LEFT, rtl=False, first=True, space_after=0)

def checklist(s, x, y, w, items, size=14.5, gap=0.04, warn=False, mark="✓", mcolor=None, tcolor=BODY):
    box, tf = tb(s, x, y, w, 5)
    mc = mcolor or (ROSE if warn else TEAL700)
    for i, it in enumerate(items):
        para(tf, [(mark + "  ", mc, True, FONT), (it, tcolor, False, FONT)],
             size=size, align=PP_ALIGN.RIGHT, first=(i == 0), space_after=int(gap * 72), line=1.12)
    return box

def card(s, x, y, w, h, emoji, title, body=None, tint='teal', items=None, title_size=15, body_size=11.5):
    rrect(s, x, y, w, h, fill=WHITE, line=LINE, radius=0.07, shadow=True)
    pad = 0.22
    icon(s, x + w - pad - 0.62, y + pad, emoji, tint, d=0.62)
    _, tf = tb(s, x + pad, y + pad + 0.78, w - 2 * pad, h - pad - 0.95)
    para(tf, [(title, INK, True, FONT)], size=title_size, align=PP_ALIGN.RIGHT, first=True, space_after=4)
    if body:
        para(tf, [(body, MUTED, False, FONT)], size=body_size, align=PP_ALIGN.RIGHT, space_after=0, line=1.18)
    if items:
        for it in items:
            para(tf, [("•  ", TINTS[tint][1], True, FONT), (it, BODY, False, FONT)],
                 size=body_size, align=PP_ALIGN.RIGHT, space_after=2, line=1.1)

def pic_row(s, paths, height, y, gap, center=PW / 2, border=False):
    pics = []
    for p in paths:
        pic = s.shapes.add_picture(p, 0, Inches(y), height=Inches(height))
        if border:
            pic.line.color.rgb = LINE; pic.line.width = Pt(1)
        pics.append(pic)
    total = sum(p.width for p in pics) + Emu(Inches(gap)) * (len(pics) - 1)
    start = Emu(Inches(center)) - total // 2
    cx = start
    for p in pics:
        p.left = int(cx)
        cx = int(cx) + p.width + int(Emu(Inches(gap)))
    return pics

def web_frame(s, path, x, y, w, caption=None):
    """Place a web screenshot inside a simple browser frame."""
    bar_h = 0.26
    pic = s.shapes.add_picture(path, Inches(x), Inches(y + bar_h), width=Inches(w))
    ph = pic.height / 914400.0
    frame = rrect(s, x - 0.04, y - 0.02, w + 0.08, bar_h + ph + 0.06, fill=WHITE, line=LINE, radius=0.04, shadow=True)
    frame._element.addprevious(frame._element)  # keep
    # move frame behind picture
    sp = frame._element
    sp.getparent().remove(sp)
    pic._element.addprevious(sp)
    bar = rrect(s, x - 0.04, y - 0.02, w + 0.08, bar_h, fill=RGBColor(0xEE, 0xF3, 0xF2), line=None, radius=0.04)
    bsp = bar._element; bsp.getparent().remove(bsp); pic._element.addprevious(bsp)
    for i, col in enumerate([RGBColor(0xFF,0x5F,0x57), RGBColor(0xFE,0xBC,0x2E), RGBColor(0x28,0xC8,0x40)]):
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.04 + i * 0.16), Inches(y + 0.06), Inches(0.1), Inches(0.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background(); dot.shadow.inherit = False
    if caption:
        _, tf = tb(s, x - 0.1, y + bar_h + ph + 0.07, w + 0.2, 0.3)
        para(tf, [(caption, MUTED, False, FONT)], size=11, align=PP_ALIGN.CENTER, first=True, space_after=0)
    return y + bar_h + ph

def caption(s, text, y):
    _, tf = tb(s, MX, y, PW - 2 * MX, 0.4)
    para(tf, [(text, MUTED, False, FONT)], size=12.5, align=PP_ALIGN.CENTER, first=True, space_after=0)

P = lambda n: os.path.join(IMG, n)

# ======================================================
# SLIDE 1 — COVER
# ======================================================
s = slide()
bgr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
gradient(bgr, TEAL500, TEAL900, 60)
# decorative circles
for (cx, cy, d, a) in [(-1.5, -1.8, 6.5, '26'), (11.5, 6.0, 5.5, '20')]:
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(cy), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = WHITE
    o.fill.transparency = 0  # placeholder
    # set alpha via xml
    sf = o.fill.fore_color._xFill
    srgb = sf.find(qn('a:srgbClr'))
    etree.SubElement(srgb, qn('a:alpha'), {'val': str(int(a) * 1000)})
    o.line.fill.background(); o.shadow.inherit = False

_, tf = tb(s, 2.0, 1.55, PW - 4.0, 0.95, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("FitLife  ", WHITE, True, FONT), ("🌿", WHITE, False, EMOJI)],
     size=46, align=PP_ALIGN.CENTER, rtl=False, first=True, space_after=0)
_, tf = tb(s, 1.2, 2.55, PW - 2.4, 1.5, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("منصة FitLife لإدارة اللياقة البدنية والتغذية", WHITE, True, FONT)],
     size=34, align=PP_ALIGN.CENTER, first=True, space_after=0, line=1.25)
_, tf = tb(s, 2.3, 4.15, PW - 4.6, 1.1)
para(tf, [("حل تقني متكامل يربط العملاء، المدربين، أخصائيي التغذية، والإدارة داخل منصة واحدة "
           "تدعم التدريب، التغذية، المتابعة، التواصل، والدفع المحلي.",
           RGBColor(0xEA, 0xFF, 0xFB), False, FONT)],
     size=16, align=PP_ALIGN.CENTER, first=True, space_after=0, line=1.5)
# chips
chips = ["مشروع تخرج — هندسة البرمجيات", "تطبيق موبايل + لوحة تحكم ويب"]
cw = 3.6; gap = 0.3; total = cw * 2 + gap; sx = (PW - total) / 2
for i, c in enumerate(chips):
    ch = rrect(s, sx + i * (cw + gap), 5.55, cw, 0.55, fill=None, line=WHITE, radius=0.5, line_w=1.25)
    ch.line.color.rgb = WHITE
    _, tf = tb(s, sx + i * (cw + gap), 5.55, cw, 0.55, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(c, WHITE, True, FONT)], size=13.5, align=PP_ALIGN.CENTER, first=True, space_after=0)

# ======================================================
# SLIDE 2 — 01 خلفية المشروع
# ======================================================
s = slide()
header(s, "01", "خلفية المشروع", [("لماذا ", INK), ("FitLife؟", TEAL700)])
_, tf = tb(s, 4.6, 1.95, PW - MX - 4.6, 1.4)
para(tf, [("يشهد المجتمع اليوم وعيًا متزايدًا بأهمية الصحة واللياقة والتغذية السليمة، لكن الحصول "
           "على متابعة احترافية ومنظمة ما زال يواجه تحديات حقيقية في السوق المحلي.", BODY)],
     size=16, align=PP_ALIGN.RIGHT, first=True, space_after=0, line=1.45)
checklist(s, 4.6, 3.5, PW - MX - 4.6, [
    "زيادة الاهتمام بالصحة ونمط الحياة المتوازن",
    "الحاجة إلى متابعة شخصية من مختصين موثوقين",
    "غياب منصة محلية موحدة تجمع التدريب والتغذية",
    "صعوبة تنظيم الخطط والقياسات والتواصل والدفع في مكان واحد",
], size=15, gap=0.12)
s.shapes.add_picture(P("mobile-home.png"), Inches(0.8), Inches(1.85), height=Inches(5.1))
footer(s)

# ======================================================
# SLIDE 3 — 02 النظام الحالي
# ======================================================
s = slide()
header(s, "02", "النظام الحالي", [("كيف تتم العملية ", INK), ("حاليًا؟", TEAL700)])
_, tf = tb(s, MX, 1.78, PW - 2 * MX, 0.5)
para(tf, [("يعتمد المستخدمون والمختصون غالبًا على أدوات تقليدية منفصلة وغير مترابطة:", BODY)],
     size=15, align=PP_ALIGN.RIGHT, first=True, space_after=0)
cur = [("💬", "تواصل مشتت", "الخطط تُرسل عبر WhatsApp أو Telegram كرسائل وملفات PDF.", 'teal'),
       ("📝", "تسجيل عشوائي", "متابعة التقدم عبر صور ورسائل دون نظام واضح.", 'amber'),
       ("📉", "لا لوحة موحّدة", "غياب لوحة لعرض القياسات والنتائج بشكل منظم.", 'rose'),
       ("💳", "دفع غير منظم", "عدم وجود نظام دفع محلي مدمج وموثوق.", 'blue'),
       ("🗂️", "سجلات ضائعة", "صعوبة الرجوع إلى السجلات والبيانات السابقة.", 'violet'),
       ("🔁", "تكرار يدوي", "إعادة العمل نفسه يدويًا لكل عميل على حدة.", 'teal')]
cw = (PW - 2 * MX - 2 * 0.3) / 3; ch = 1.95; gx = 0.3; gy = 0.3; y0 = 2.4
for i, (e, t, b, tn) in enumerate(cur):
    col = i % 3; row = i // 3
    x = PW - MX - cw - col * (cw + gx)
    card(s, x, y0 + row * (ch + gy), cw, ch, e, t, b, tn)
footer(s)

# ======================================================
# SLIDE 4 — 03 مشكلة المشروع
# ======================================================
s = slide()
header(s, "03", "تعريف المشكلة", [("المشكلة ", INK), ("الأساسية", TEAL700)])
hb = rrect(s, MX, 1.95, PW - 2 * MX, 1.25, fill=RGBColor(0xE9, 0xFB, 0xF6), line=RGBColor(0x99, 0xF6, 0xE4), radius=0.08)
_, tf = tb(s, MX + 0.4, 1.95, PW - 2 * MX - 0.8, 1.25, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("لا توجد منصة رقمية موحدة تربط بين العميل والمدرب وأخصائي التغذية داخل بيئة واحدة "
           "تدعم إدارة الخطط، تتبع التقدم، التواصل، والدفع المحلي.", TEAL900, True)],
     size=18, align=PP_ALIGN.RIGHT, first=True, space_after=0, line=1.4)
colw = (PW - 2 * MX - 0.5) / 2
checklist(s, PW - MX - colw, 3.55, colw, [
    "تشتت وسائل التواصل بين أطراف متعددة",
    "عدم وجود نظام موحد لإدارة الخطط",
    "ضعف القدرة على تتبع تقدم العميل",
    "غياب الرسوم البيانية والمؤشرات الصحية"], size=15.5, gap=0.18, warn=True, mark="!")
checklist(s, MX, 3.55, colw, [
    "صعوبة الوصول إلى مختصين موثوقين",
    "غياب آلية دفع محلية مناسبة",
    "عدم وجود تنبيهات وتذكيرات فعّالة",
    "صعوبة تحليل وتقييم النتائج"], size=15.5, gap=0.18, warn=True, mark="!")
footer(s)

# ======================================================
# SLIDE 5 — 04 الحل المقترح
# ======================================================
s = slide()
header(s, "04", "الحل المقترح", [("منصة ", INK), ("FitLife", TEAL700)])
_, tf = tb(s, 4.7, 1.95, PW - MX - 4.7, 1.3)
para(tf, [("منصة رقمية متكاملة لإدارة خدمات اللياقة والتغذية عبر ", BODY),
          ("تطبيق موبايل للعملاء", INK, True), (" و", BODY),
          ("لوحة تحكم ويب", INK, True), (" للمختصين والإدارة.", BODY)],
     size=16, align=PP_ALIGN.RIGHT, first=True, space_after=0, line=1.45)
pills = ["🏋️ خطط تدريبية", "🥗 خطط غذائية", "📈 تتبع القياسات",
         "💬 محادثة مباشرة", "🔔 إشعارات وتنبيهات", "💳 دفع محلي (معاملات)"]
pw = (PW - MX - 4.7 - 0.25) / 2; ph = 0.6; y0 = 3.35
for i, p in enumerate(pills):
    col = i % 2; row = i // 2
    x = PW - MX - pw - col * (pw + 0.25)
    e, txt = p.split(" ", 1)
    pr = rrect(s, x, y0 + row * (ph + 0.18), pw, ph, fill=MINT100, line=None, radius=0.5)
    _, tf = tb(s, x + 0.2, y0 + row * (ph + 0.18), pw - 0.4, ph, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(e + " ", TEAL700, False, EMOJI), (txt, TEAL800, True, FONT)],
         size=13.5, align=PP_ALIGN.RIGHT, first=True, space_after=0)
pic_row(s, [P("mobile-workouts.png"), P("mobile-nutrition-plan.png")], 4.9, 1.85, 0.25, center=2.55)
footer(s)

# ======================================================
# SLIDE 6 — 05 القيمة المقترحة
# ======================================================
s = slide()
header(s, "05", "القيمة المقترحة", [("ما الذي يميّز ", INK), ("FitLife؟", TEAL700)])
val = [("🙋", "للعميل", 'teal', ["وصول إلى مختصين موثوقين", "خطط تدريب وغذاء منظمة",
                                   "متابعة الإنجاز اليومي", "تقدم عبر مؤشرات ورسوم بيانية"]),
       ("🧑‍🏫", "للمختص", 'blue', ["إدارة العملاء من مكان واحد", "إنشاء باقات وبرامج متكاملة",
                                    "متابعة التزام العملاء", "إدارة الأرباح وطلبات السحب"]),
       ("🌍", "للسوق المحلي", 'amber', ["دعم الدفع المحلي", "رقمنة قطاع التدريب والتغذية",
                                         "تقليل الرسائل العشوائية", "تجربة موحّدة واحترافية"])]
cw = (PW - 2 * MX - 2 * 0.35) / 3; ch = 3.7; y0 = 2.15
for i, (e, t, tn, items) in enumerate(val):
    x = PW - MX - cw - i * (cw + 0.35)
    card(s, x, y0, cw, ch, e, t, None, tn, items=items, title_size=18, body_size=13)
footer(s)

# ======================================================
# SLIDE 7 — 06 المستخدمون
# ======================================================
s = slide()
header(s, "06", "المستخدمون المستهدفون", [("أربعة ", INK), ("أدوار", TEAL700), (" داخل النظام", INK)])
roles = [("🙋", "العميل", 'teal', "الاشتراك، عرض الخطط اليومية، تسجيل التمارين والوجبات، رفع القياسات والصور، ومتابعة المؤشرات."),
         ("🏋️", "المدرب", 'blue', "إنشاء مكتبة تمارين، بناء برامج تدريبية، جدولة التمارين، ومتابعة أداء العميل."),
         ("🥗", "أخصائي التغذية", 'amber', "إنشاء مكتبة أغذية ووجبات، بناء برامج غذائية، ومراجعة السعرات والماكروز اليومية."),
         ("🛡️", "المسؤول", 'violet', "إدارة المستخدمين، مراقبة المدفوعات والاشتراكات، إدارة المحتوى، ومتابعة طلبات السحب.")]
cw = (PW - 2 * MX - 3 * 0.3) / 4; ch = 3.3; y0 = 2.15
for i, (e, t, tn, b) in enumerate(roles):
    x = PW - MX - cw - i * (cw + 0.3)
    card(s, x, y0, cw, ch, e, t, b, tn, title_size=16, body_size=12.5)
_, tf = tb(s, MX, 5.7, PW - 2 * MX, 0.5)
para(tf, [("تصميم قائم على الأدوار (Role-Based) يمنح كل مستخدم صلاحيات ووظائف محددة — تنظيم أعلى، أمان أكبر، وقابلية للتوسع.", MUTED)],
     size=12.5, align=PP_ALIGN.CENTER, first=True, space_after=0)
footer(s)

# ======================================================
# SLIDE 8 — 07 الوحدات
# ======================================================
s = slide()
header(s, "07", "وظائف النظام", [("الوحدات ", INK), ("الرئيسية", TEAL700)])
mods = ["🔐 الحسابات والمصادقة", "🛒 سوق المدربين والمختصين", "💳 الاشتراكات والدفع",
        "🏋️ البرامج التدريبية", "🥗 البرامج الغذائية", "📈 القياسات والتقدم",
        "💬 المحادثات والإشعارات", "🖥️ لوحة تحكم الإدارة", "👛 المحافظ وطلبات السحب",
        "📚 مكتبات التمارين والأغذية"]
colw = 2.9; ph = 0.55; y0 = 2.1; xL = 6.95
for i, m in enumerate(mods):
    col = i % 2; row = i // 2
    x = xL + (colw + 0.25) - col * (colw + 0.25)
    e, txt = m.split(" ", 1)
    rrect(s, x, y0 + row * (ph + 0.16), colw, ph, fill=MINT100, line=None, radius=0.4)
    _, tf = tb(s, x + 0.18, y0 + row * (ph + 0.16), colw - 0.36, ph, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(e + " ", TEAL700, False, EMOJI), (txt, TEAL800, True, FONT)],
         size=12.5, align=PP_ALIGN.RIGHT, first=True, space_after=0)
web_frame(s, P("web-coach-dashboard.png"), 0.7, 2.25, 5.7, "لوحة تحكم المدرب — نظرة عامة على النشاط والعملاء")
footer(s)

# ======================================================
# SLIDE 9 — 08 رحلة العميل
# ======================================================
s = slide()
header(s, "08", "رحلة العميل", [("من ", INK), ("الهدف", TEAL700), (" إلى ", INK), ("النتيجة", TEAL700)])
steps = ["إنشاء حساب وتعبئة البيانات", "تصفح المدربين والمختصين", "اختيار المختص والباقة",
         "إتمام الدفع محليًا", "استلام الخطة المناسبة", "تسجيل التمارين والوجبات",
         "رفع القياسات وصور التقدم", "متابعة النتائج بالرسوم البيانية", "التواصل مع المختص عند الحاجة"]
colw = 3.55; sh = 0.78; y0 = 2.1; xL = 4.55
for i, st in enumerate(steps):
    col = i % 2; row = i // 2
    x = xL + (colw + 0.2) - col * (colw + 0.2)
    rrect(s, x, y0 + row * (sh + 0.12), colw, sh, fill=WHITE, line=LINE, radius=0.12, shadow=True)
    nb = rrect(s, x + colw - 0.62, y0 + row * (sh + 0.12) + 0.16, 0.46, 0.46, radius=0.25)
    gradient(nb, TEAL500, TEAL800, 45)
    _, tf = tb(s, x + colw - 0.62, y0 + row * (sh + 0.12) + 0.16, 0.46, 0.46, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(str(i + 1), WHITE, True, FONT)], size=15, align=PP_ALIGN.CENTER, rtl=False, first=True, space_after=0)
    _, tf = tb(s, x + 0.18, y0 + row * (sh + 0.12), colw - 0.85, sh, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(st, BODY, False)], size=13, align=PP_ALIGN.RIGHT, first=True, space_after=0, line=1.1)
s.shapes.add_picture(P("mobile-onboarding.png"), Inches(0.75), Inches(2.0), height=Inches(4.9))
footer(s)

# ======================================================
# SLIDE 10 — 09 التقنيات
# ======================================================
s = slide()
header(s, "09", "البنية التقنية", [("التقنيات ", INK), ("المستخدمة", TEAL700)])
tech = [("⚙️", "الخلفية (Backend)", ["PHP", "Laravel", "MySQL"], 'teal'),
        ("🖥️", "لوحة التحكم", ["Vue.js", "TypeScript"], 'blue'),
        ("📱", "تطبيق الموبايل", ["Flutter", "Dart"], 'amber'),
        ("☁️", "خدمات داعمة", ["Firebase", "Moamalat"], 'violet')]
cw = (PW - 2 * MX - 3 * 0.3) / 4; ch = 2.6; y0 = 2.15
for i, (e, t, tags, tn) in enumerate(tech):
    x = PW - MX - cw - i * (cw + 0.3)
    rrect(s, x, y0, cw, ch, fill=WHITE, line=LINE, radius=0.08, shadow=True)
    icon(s, x + cw - 0.22 - 0.6, y0 + 0.22, e, tn, d=0.6)
    _, tf = tb(s, x + 0.22, y0 + 1.0, cw - 0.44, 1.5)
    para(tf, [(t, INK, True)], size=14, align=PP_ALIGN.RIGHT, first=True, space_after=6)
    for tg in tags:
        para(tf, [(tg, TEAL700, True)], size=13.5, align=PP_ALIGN.RIGHT, space_after=3)
wb = rrect(s, MX, 5.0, PW - 2 * MX, 1.5, fill=WHITE, line=LINE, radius=0.07, shadow=True)
icon(s, PW - MX - 0.22 - 0.6, 5.22, "🛠️", 'teal', d=0.6)
_, tf = tb(s, MX + 0.3, 5.18, PW - 2 * MX - 1.1, 0.5)
para(tf, [("أدوات التطوير", INK, True)], size=14, align=PP_ALIGN.RIGHT, first=True, space_after=4)
_, tf = tb(s, MX + 0.3, 5.7, PW - 2 * MX - 0.6, 0.6)
para(tf, [("VS Code   •   Git & GitHub   •   Postman   •   Draw.io   •   MockFlow", BODY)],
     size=14, align=PP_ALIGN.RIGHT, first=True, space_after=0)
footer(s)

# ======================================================
# SLIDE 11 — 10 لماذا هذه التقنيات (جدول)
# ======================================================
s = slide()
header(s, "10", "القرارات التقنية", [("لماذا ", INK), ("هذه التقنيات؟", TEAL700)])
rows = [("المكوّن", "التقنية", "سبب الاستخدام"),
        ("الواجهة الخلفية", "Laravel / PHP", "بناء API، المصادقة، منطق النظام، وقاعدة البيانات"),
        ("لوحة التحكم", "Vue.js / TypeScript", "واجهات تفاعلية وسهلة الصيانة"),
        ("تطبيق الهاتف", "Flutter / Dart", "تطبيق متعدد المنصات بأداء جيد"),
        ("قاعدة البيانات", "MySQL", "تخزين منظم للبيانات العلائقية"),
        ("الخدمات الداعمة", "Firebase", "الإشعارات وبعض خدمات التطبيق"),
        ("الدفع المحلي", "Moamalat", "تمكين المستخدم من الدفع بطرق محلية")]
tw = PW - 2 * MX; tx = MX; ty = 2.1
gtab = s.shapes.add_table(len(rows), 3, Inches(tx), Inches(ty), Inches(tw), Inches(4.2)).table
gtab.columns[0].width = Inches(2.6)
gtab.columns[1].width = Inches(3.2)
gtab.columns[2].width = Inches(tw - 5.8)
# disable default style banding by setting first_row false later; fill manually
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = gtab.cell(ri, ci)
        cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.12)
        cell.margin_top = Inches(0.06); cell.margin_bottom = Inches(0.06)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = TEAL700
            col = WHITE; bold = True; size = 14
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else RGBColor(0xF2, 0xFA, 0xF8)
            if ci == 0:
                col = INK; bold = True; size = 13.5
            elif ci == 1:
                col = TEAL700; bold = True; size = 13.5
            else:
                col = BODY; bold = False; size = 13
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT; set_rtl(p)
        r = p.add_run(); r.text = val
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONT; r.font.color.rgb = col
        _cs(r, FONT)
footer(s)

# ======================================================
# SLIDE 12 — 11 المعمارية
# ======================================================
s = slide()
header(s, "11", "معمارية النظام", [("كيف يعمل ", INK), ("FitLife؟", TEAL700)])

def node(x, y, w, h, title, desc=None, kind='soft'):
    if kind == 'core':
        sh = rrect(s, x, y, w, h, radius=0.1); gradient(sh, TEAL500, TEAL800, 45)
        tc, dc = WHITE, RGBColor(0xDF, 0xF7, 0xF3)
    elif kind == 'client':
        sh = rrect(s, x, y, w, h, fill=RGBColor(0xE9, 0xFB, 0xF6), line=RGBColor(0x99, 0xF6, 0xE4), radius=0.1)
        tc, dc = INK, MUTED
    else:
        sh = rrect(s, x, y, w, h, fill=WHITE, line=LINE, radius=0.1, shadow=True)
        tc, dc = INK, MUTED
    _, tf = tb(s, x + 0.2, y, w - 0.4, h, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(title, tc, True)], size=16, align=PP_ALIGN.CENTER, first=True, space_after=(3 if desc else 0))
    if desc:
        para(tf, [(desc, dc, False)], size=11.5, align=PP_ALIGN.CENTER, space_after=0, line=1.15)

def arrow_down(x, y):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(0.32), Inches(0.34))
    a.fill.solid(); a.fill.fore_color.rgb = TEAL600; a.line.fill.background(); a.shadow.inherit = False

nw = 4.7
node(PW - MX - nw, 2.05, nw, 1.05, "📱 تطبيق الموبايل", "العملاء — الخطط، التسجيل، التقدم، المحادثة", 'client')
node(MX, 2.05, nw, 1.05, "🖥️ لوحة تحكم الويب", "الإدارة والمدربون وأخصائيو التغذية", 'client')
arrow_down(PW - MX - nw / 2 - 0.16, 3.2)
arrow_down(MX + nw / 2 - 0.16, 3.2)
node((PW - 8.2) / 2, 3.7, 8.2, 1.1, "🔗 الواجهة الخلفية — Backend API (Laravel)",
     "المصادقة · منطق النظام · الدفع · إدارة البيانات · التكامل مع Firebase & Moamalat", 'core')
arrow_down(PW / 2 - 0.16, 4.9)
node((PW - 4.6) / 2, 5.4, 4.6, 0.85, "🗄️ قاعدة البيانات — MySQL", None, 'client')
_, tf = tb(s, MX, 6.45, PW - 2 * MX, 0.4)
para(tf, [("فصل واضح للمكونات: كل جزء له مسؤولية محددة، والتواصل بينها يتم عبر API.", MUTED)],
     size=12.5, align=PP_ALIGN.CENTER, first=True, space_after=0)
footer(s)

# ======================================================
# SLIDE 13 — 12 واجهات الموبايل
# ======================================================
s = slide()
header(s, "12", "الواجهات وتجربة المستخدم", [("تطبيق ", INK), ("العميل", TEAL700)])
pic_row(s, [P("mobile-workouts.png"), P("mobile-exercise-detail.png"),
            P("mobile-nutrition-plan.png"), P("mobile-chat.png")], 4.55, 1.85, 0.3)
caption(s, "جدول التمارين · تفاصيل التمرين · الخطة الغذائية والسعرات · المحادثة مع المختص — تصميم بسيط يدعم العربية بالكامل.", 6.6)
footer(s)

# ======================================================
# SLIDE 14 — 12 لوحة التحكم
# ======================================================
s = slide()
header(s, "12", "الواجهات وتجربة المستخدم", [("لوحة تحكم ", INK), ("المختصين والإدارة", TEAL700)])
w = 5.55
web_frame(s, P("web-exercise-library.png"), PW - MX - w, 1.95, w, "مكتبة التمارين")
web_frame(s, P("web-recipes.png"), MX, 1.95, w, "مكتبة الوصفات والوجبات")
web_frame(s, P("web-client-details.png"), PW - MX - w, 4.55, w, "ملف العميل والقياسات والأهداف")
web_frame(s, P("web-packages.png"), MX, 4.55, w, "إدارة الخدمات والباقات")
footer(s)

# ======================================================
# SLIDE 15 — 13 المتطلبات غير الوظيفية
# ======================================================
s = slide()
header(s, "13", "معايير الجودة", [("المتطلبات ", INK), ("غير الوظيفية", TEAL700)])
nf = [("⚡", "الأداء", "استجابة سريعة للعمليات الأساسية.", 'teal'),
      ("🔒", "الأمان والخصوصية", "حماية البيانات الصحية، تشفير كلمات المرور، ومنع الوصول غير المصرح به.", 'rose'),
      ("🎯", "سهولة الاستخدام", "واجهات واضحة ومتناسقة بدعم كامل للغة العربية.", 'blue'),
      ("🔗", "التوافقية", "العمل على الهواتف والمتصفحات الحديثة.", 'amber'),
      ("🛡️", "الموثوقية", "الحفاظ على التوفر مع دعم النسخ الاحتياطي.", 'violet'),
      ("📈", "قابلية التوسع", "استيعاب زيادة المستخدمين والبيانات مستقبلًا.", 'teal')]
cw = (PW - 2 * MX - 2 * 0.3) / 3; ch = 1.95; y0 = 2.2
for i, (e, t, b, tn) in enumerate(nf):
    col = i % 3; row = i // 3
    x = PW - MX - cw - col * (cw + 0.3)
    card(s, x, y0 + row * (ch + 0.3), cw, ch, e, t, b, tn)
footer(s)

# ======================================================
# SLIDE 16 — 14 الجدوى
# ======================================================
s = slide()
header(s, "14", "دراسة الجدوى", [("لماذا يمكن أن ", INK), ("ينجح FitLife؟", TEAL700)])
fz = [("💰", "الجدوى الاقتصادية", 'amber', ["طلب متزايد على خدمات اللياقة", "بيع المختصين خدماتهم مباشرة",
                                              "دعم الدفع المحلي", "نموذج عمل قابل للنمو"]),
      ("⚙️", "الجدوى التشغيلية", 'teal', ["توحيد العمليات في منصة واحدة", "تقليل الفوضى بين التطبيقات",
                                            "تسهيل إدارة العملاء والخطط", "إدارة الاشتراكات بسلاسة"]),
      ("🧩", "الجدوى الفنية", 'blue', ["تقنيات معروفة ومستقرة", "دعم الويب والموبايل",
                                        "قابلية النظام للتطوير", "بنية واضحة قابلة للصيانة"])]
cw = (PW - 2 * MX - 2 * 0.35) / 3; ch = 3.7; y0 = 2.15
for i, (e, t, tn, items) in enumerate(fz):
    x = PW - MX - cw - i * (cw + 0.35)
    card(s, x, y0, cw, ch, e, t, None, tn, items=items, title_size=17, body_size=13)
footer(s)

# ======================================================
# SLIDE 17 — 15 الميزة التنافسية
# ======================================================
s = slide()
header(s, "15", "الميزة التنافسية", [("ما الذي يجعل FitLife ", INK), ("مختلفًا؟", TEAL700)])
colw = (PW - 2 * MX - 0.6) / 2
checklist(s, PW - MX - colw, 2.4, colw, [
    "مصمّم ليناسب السوق المحلي",
    "يدعم الدفع المحلي عبر معاملات",
    "يجمع التدريب والتغذية والمتابعة والتواصل في منصة واحدة",
    "يخدم العميل والمختص والإدارة معًا"], size=17, gap=0.32)
checklist(s, MX, 2.4, colw, [
    "يحوّل المتابعة من رسائل عشوائية إلى نظام منظم",
    "يعتمد على البيانات والرسوم البيانية لقياس التقدم",
    "قابل للتوسع بإضافة مزايا ذكية مستقبلًا",
    "تجربة عربية كاملة وحديثة"], size=17, gap=0.32)
footer(s)

# ======================================================
# SLIDE 18 — 16 المنهجية
# ======================================================
s = slide()
header(s, "16", "منهجية التطوير", [("نموذج الشلال ", INK), ("المعدّل", TEAL700)])
phases = ["تحليل المتطلبات", "تصميم النظام", "التنفيذ", "الاختبار", "النشر", "الصيانة والتحسين"]
n = len(phases); gap = 0.2; cw = (PW - 2 * MX - gap * (n - 1)) / n; y0 = 2.15; chh = 1.5
for i, ph in enumerate(phases):
    x = PW - MX - cw - i * (cw + gap)
    rrect(s, x, y0, cw, chh, fill=WHITE, line=LINE, radius=0.1, shadow=True)
    nb = rrect(s, x + cw / 2 - 0.27, y0 + 0.22, 0.54, 0.54, radius=0.25); gradient(nb, TEAL500, TEAL800, 45)
    _, tf = tb(s, x + cw / 2 - 0.27, y0 + 0.22, 0.54, 0.54, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(str(i + 1), WHITE, True)], size=16, align=PP_ALIGN.CENTER, rtl=False, first=True, space_after=0)
    _, tf = tb(s, x + 0.08, y0 + 0.85, cw - 0.16, 0.6, anchor=MSO_ANCHOR.TOP)
    para(tf, [(ph, BODY, True)], size=12.5, align=PP_ALIGN.CENTER, first=True, space_after=0, line=1.05)
box = rrect(s, MX, 4.1, PW - 2 * MX, 2.0, fill=RGBColor(0xE9, 0xFB, 0xF6), line=RGBColor(0x99, 0xF6, 0xE4), radius=0.07)
_, tf = tb(s, MX + 0.4, 4.35, PW - 2 * MX - 0.8, 1.5)
para(tf, [("لماذا هذا النموذج؟", TEAL900, True)], size=17, align=PP_ALIGN.RIGHT, first=True, space_after=8)
para(tf, [("متطلبات المشروع كانت واضحة من البداية، والنموذج مناسب لمشاريع التخرج، يوفر تسلسلًا منظمًا "
           "للمراحل، ويسمح بالرجوع المحدود للمراحل السابقة عند الحاجة للتصحيح.", TEAL900)],
     size=15, align=PP_ALIGN.RIGHT, space_after=0, line=1.4)
footer(s)

# ======================================================
# SLIDE 19 — 17 التطوير المستقبلي
# ======================================================
s = slide()
header(s, "17", "رؤية مستقبلية", [("التطوير ", INK), ("المستقبلي", TEAL700)])
fut = [("🤖", "توصيات ذكية", "نظام توصيات بالذكاء الاصطناعي للتمارين والوجبات.", 'violet'),
       ("⏰", "تنبيهات ذكية", "تذكيرات تكيّفية حسب سلوك المستخدم.", 'teal'),
       ("📊", "تحليل متقدم", "تحليل أعمق للتقدم الصحي وتقارير للإدارة.", 'blue'),
       ("⭐", "تقييمات ومراجعات", "تقييم المدربين وأخصائيي التغذية.", 'amber'),
       ("⌚", "الأجهزة القابلة للارتداء", "ربط النظام بالساعات وأجهزة التتبع.", 'rose'),
       ("💳", "طرق دفع إضافية", "توسيع خيارات الدفع المحلية والعالمية.", 'teal')]
cw = (PW - 2 * MX - 2 * 0.3) / 3; ch = 1.95; y0 = 2.2
for i, (e, t, b, tn) in enumerate(fut):
    col = i % 3; row = i // 3
    x = PW - MX - cw - col * (cw + 0.3)
    card(s, x, y0 + row * (ch + 0.3), cw, ch, e, t, b, tn)
footer(s)

# ======================================================
# SLIDE 20 — 18 الخاتمة
# ======================================================
s = slide()
bgr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
gradient(bgr, TEAL500, TEAL900, 60)
header(s, "18", "الخاتمة", [("FitLife كحل تقني ", WHITE), ("متكامل", RGBColor(0xBF, 0xF7, 0xEC))], light=True)
_, tf = tb(s, MX, 1.95, PW - 2 * MX, 0.8)
para(tf, [("منصة رقمية متكاملة تربط العملاء والمدربين وأخصائيي التغذية والإدارة داخل نظام واحد منظم.", WHITE)],
     size=17, align=PP_ALIGN.CENTER, first=True, space_after=0, line=1.4)
colw = (PW - 2 * MX - 0.6) / 2
checklist(s, PW - MX - colw, 2.95, colw, [
    "يعالج مشاكل النظام التقليدي", "ينظّم الخطط والمتابعة",
    "يحسّن التواصل بين العميل والمختص", "يدعم الدفع المحلي"],
    size=15.5, gap=0.22, mcolor=WHITE, tcolor=WHITE)
checklist(s, MX, 2.95, colw, [
    "يوفّر تجربة استخدام واضحة", "قابل للتوسع والتطوير مستقبلًا",
    "يساهم في التحول الرقمي للصحة والرياضة", "حل واقعي لمشكلة قائمة في السوق"],
    size=15.5, gap=0.22, mcolor=WHITE, tcolor=WHITE)
_, tf = tb(s, MX, 5.55, PW - 2 * MX, 0.9, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("شكرًا لاستماعكم — أسئلتكم ومناقشتكم محل ترحيب ", WHITE, True), ("🌿", WHITE, False, EMOJI)],
     size=24, align=PP_ALIGN.CENTER, first=True, space_after=0)

out = os.path.join(BASE, "FitLife_Presentation.pptx")
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
